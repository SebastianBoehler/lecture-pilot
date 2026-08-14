from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def pptx_supplemental_blocks(path: Path) -> list[dict]:
    presentation = Presentation(path)
    blocks = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            blocks.append(
                {
                    "kind": "paragraph",
                    "text": notes,
                    "locator": {"slide": slide_number},
                    "extraction": "native",
                }
            )
        blocks.extend(_slide_links(slide, slide_number=slide_number))
    return blocks


def _slide_links(slide, *, slide_number: int) -> list[dict]:
    blocks = []
    for shape in slide.shapes:
        if address := shape.click_action.hyperlink.address:
            blocks.append(
                _link_block(
                    address,
                    label=(shape.text.strip() if shape.has_text_frame else ""),
                    slide_number=slide_number,
                )
            )
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if address := run.hyperlink.address:
                    blocks.append(
                        _link_block(
                            address,
                            label=run.text.strip(),
                            slide_number=slide_number,
                        )
                    )
    return list({(block["url"], block["text"]): block for block in blocks}.values())


def _link_block(address: str, *, label: str, slide_number: int) -> dict:
    return {
        "kind": "link",
        "text": label or address,
        "url": address,
        "locator": {"slide": slide_number},
        "extraction": "native",
    }
