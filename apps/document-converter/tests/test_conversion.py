from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from lecturepilot_converter.conversion import convert_document


SOURCE_SHA256 = "a" * 64


def test_docx_conversion_preserves_heading_and_paragraph_order(tmp_path: Path) -> None:
    source = tmp_path / "overview.docx"
    document = Document()
    document.add_heading("Course overview", level=1)
    document.add_paragraph("Evidence-based learning.")
    document.save(source)

    manifest = convert_document(
        source,
        source_path="week-01/overview.docx",
        source_sha256=SOURCE_SHA256,
        output_root=tmp_path / "normalized",
    )

    assert manifest["source_path"] == "week-01/overview.docx"
    assert manifest["source_sha256"] == SOURCE_SHA256
    assert [(block["kind"], block["text"]) for block in manifest["blocks"]] == [
        ("heading", "Course overview"),
        ("paragraph", "Evidence-based learning."),
    ]
    revision_root = tmp_path / "normalized" / SOURCE_SHA256
    assert (revision_root / "manifest.json").is_file()
    assert (revision_root / "content.md").read_text(encoding="utf-8") == (
        "## Course overview\n\nEvidence-based learning."
    )


def test_pptx_conversion_preserves_slide_notes_and_shape_links(tmp_path: Path) -> None:
    source = tmp_path / "lecture.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Lecture one"
    body = slide.placeholders[1].text_frame
    body.text = "Bayesian evidence"
    body.paragraphs[0].runs[0].hyperlink.address = "https://example.edu/evidence"
    slide.shapes.title.click_action.hyperlink.address = "https://example.edu/course"
    slide.notes_slide.notes_text_frame.text = "Explain the prior before the likelihood."
    presentation.save(source)

    manifest = convert_document(
        source,
        source_path="slides/lecture.pptx",
        source_sha256=SOURCE_SHA256,
        output_root=tmp_path / "normalized",
    )

    assert all(block["locator"]["slide"] == 1 for block in manifest["blocks"])
    assert any(
        block.get("text") == "Explain the prior before the likelihood."
        for block in manifest["blocks"]
    )
    assert any(block.get("url") == "https://example.edu/evidence" for block in manifest["blocks"])
    assert any(block.get("url") == "https://example.edu/course" for block in manifest["blocks"])
    rendered = tmp_path / "normalized" / SOURCE_SHA256 / "rendered.pdf"
    assert rendered.read_bytes().startswith(b"%PDF-")


def test_xlsx_conversion_keeps_formulas_separate_from_cell_values(tmp_path: Path) -> None:
    source = tmp_path / "results.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Regression"
    sheet.append(["Method", "Score"])
    sheet.append(["OLS", 0.9])
    sheet["B3"] = "=SUM(B2)"
    workbook.save(source)

    manifest = convert_document(
        source,
        source_path="data/results.xlsx",
        source_sha256=SOURCE_SHA256,
        output_root=tmp_path / "normalized",
    )

    table = next(block for block in manifest["blocks"] if block["kind"] == "table")
    assert table["locator"] == {"sheet": "Regression", "cell_range": "A1:B3"}
    formula = next(cell for cell in table["cells"] if (cell["row"], cell["column"]) == (3, 2))
    assert formula == {"row": 3, "column": 2, "value": None, "formula": "=SUM(B2)"}


def test_image_only_pptx_uses_selective_ocr(monkeypatch, tmp_path: Path) -> None:
    image = tmp_path / "scan.png"
    Image.new("RGB", (1200, 800), "white").save(image)
    source = tmp_path / "scanned-slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image), 0, 0, width=Inches(10))
    presentation.save(source)
    monkeypatch.setenv("LECTUREPILOT_OCR_URL", "http://ocr:8080")
    monkeypatch.setattr(
        "lecturepilot_converter.ocr_client.PaddleOcrClient.extract_markdown",
        lambda _self, _image: "## Gescannte Folie\n\nBayes-Regel",
    )

    manifest = convert_document(
        source,
        source_path="slides/scanned-slides.pptx",
        source_sha256=SOURCE_SHA256,
        output_root=tmp_path / "normalized",
    )

    ocr = next(block for block in manifest["blocks"] if block["extraction"] == "ocr")
    assert ocr["locator"]["slide"] == 1
    assert "Bayes-Regel" in ocr["text"]
